"""PowerPoint export: プロジェクトデータから診断レポートPPTを生成"""
import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .templates import AREAS, AREA_ORDER, rank_from_score
from .scoring import score_all

# Palette
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
INDIGO_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
EMERALD_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)

HEADER = "Arial"
BODY = "Arial"


def _rect(slide, x, y, w, h, fill=None, line=None, corner=None):
    if corner:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.adjustments[0] = corner
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, text, *, size=14, bold=False, color=SLATE_800,
          font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _title_bar(slide, title_text, num=None):
    _rect(slide, Inches(0), Inches(0), Inches(5.5), Inches(0.7), fill=SLATE_800)
    if num:
        _text(slide, Inches(0.3), Inches(0), Inches(0.6), Inches(0.7),
              f"#{num}", size=13, bold=True, color=INDIGO, font=HEADER,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(0.9), Inches(0), Inches(4.5), Inches(0.7),
              title_text, size=17, bold=True, color=WHITE, font=HEADER,
              anchor=MSO_ANCHOR.MIDDLE)
    else:
        _text(slide, Inches(0.3), Inches(0), Inches(5.1), Inches(0.7),
              title_text, size=17, bold=True, color=WHITE, font=HEADER,
              anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(5.5), Inches(0.58), Inches(7.8), Inches(0.04), fill=SLATE_200)


def _footer(slide, num, total):
    _text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.35),
          "AIO/SEO 現状診断レポート", size=9, color=SLATE_500, italic=True)
    _text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.35),
          f"{num:02d} / {total:02d}", size=9, color=SLATE_500, align=PP_ALIGN.RIGHT)


def build_pptx(project: dict) -> bytes:
    """プロジェクトdictから.pptxバイト列を生成"""
    sc = score_all(project)
    rank, rank_desc, rank_color_hex = rank_from_score(sc["total"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    TOTAL = 8  # 8スライド構成の軽量版

    client = project.get("client_name", "{クライアント名}")
    date_str = datetime.now().strftime("%Y年%m月%d日")

    # -------- Slide 1: Cover --------
    s1 = prs.slides.add_slide(blank)
    _rect(s1, Inches(0), Inches(0), SW, SH, fill=SLATE_900)
    _text(s1, Inches(10.5), Inches(0.2), Inches(2.8), Inches(2.5),
          "01", size=220, bold=True, color=SLATE_800, font=HEADER,
          align=PP_ALIGN.RIGHT)
    _rect(s1, Inches(0.75), Inches(2.1), Inches(0.15), Inches(1.4), fill=INDIGO)
    _text(s1, Inches(1.1), Inches(2.1), Inches(11), Inches(0.5),
          "AIO / SEO Diagnostics Report", size=14, bold=True,
          color=INDIGO, font=HEADER)
    _text(s1, Inches(1.1), Inches(2.6), Inches(11), Inches(1.4),
          f"{client} 様", size=40, bold=True, color=WHITE, font=HEADER)
    _text(s1, Inches(1.1), Inches(3.7), Inches(11), Inches(1.2),
          "AIO/SEO 現状診断レポート", size=44, bold=True, color=WHITE, font=HEADER)
    _text(s1, Inches(1.1), Inches(5.1), Inches(11), Inches(0.5),
          "Webサイトリニューアル提案の前提として", size=16, color=SLATE_300,
          font=BODY, italic=True)
    _rect(s1, Inches(0), Inches(6.7), SW, Inches(0.8), fill=SLATE_800)
    _text(s1, Inches(0.75), Inches(6.7), Inches(8), Inches(0.8),
          date_str, size=13, color=WHITE, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    _text(s1, Inches(9), Inches(6.7), Inches(3.5), Inches(0.8),
          "Goodpatch Inc.", size=13, bold=True, color=WHITE, font=HEADER,
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # -------- Slide 2: Executive Summary --------
    s2 = prs.slides.add_slide(blank)
    _rect(s2, Inches(0), Inches(0), SW, SH, fill=SLATE_50)
    _title_bar(s2, "エグゼクティブサマリー", num=2)

    # Score panel
    _rect(s2, Inches(0.5), Inches(1.2), Inches(5.5), Inches(2.9),
          fill=SLATE_800, corner=0.04)
    _text(s2, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
          "TOTAL SCORE", size=11, bold=True, color=INDIGO, font=HEADER,
          align=PP_ALIGN.CENTER)
    _text(s2, Inches(0.5), Inches(1.7), Inches(5.5), Inches(1.5),
          f"{sc['total']:.0f}", size=96, bold=True, color=WHITE, font=HEADER,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s2, Inches(0.5), Inches(3.15), Inches(5.5), Inches(0.4),
          "/ 100 点", size=14, color=SLATE_300, align=PP_ALIGN.CENTER)
    _rect(s2, Inches(2.35), Inches(3.55), Inches(1.8), Inches(0.5),
          fill=INDIGO, corner=0.5)
    _text(s2, Inches(2.35), Inches(3.55), Inches(1.8), Inches(0.5),
          f"Rank {rank}", size=14, bold=True, color=WHITE, font=HEADER,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right: area scores
    _text(s2, Inches(6.4), Inches(1.2), Inches(6.4), Inches(0.4),
          "領域別スコア", size=13, bold=True, color=SLATE_800, font=HEADER)
    for i, key in enumerate(AREA_ORDER):
        a = AREAS[key]
        row_y = Inches(1.7 + i * 0.42)
        _rect(s2, Inches(6.4), row_y, Inches(6.4), Inches(0.38),
              fill=WHITE if i % 2 == 0 else SLATE_100)
        color_rgb = RGBColor.from_string(a["color"].replace("#", ""))
        _rect(s2, Inches(6.4), row_y, Inches(0.08), Inches(0.38), fill=color_rgb)
        _text(s2, Inches(6.6), row_y, Inches(3.2), Inches(0.38),
              f"領域{a['num']}: {a['name']}", size=11, color=SLATE_800,
              font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        _text(s2, Inches(9.8), row_y, Inches(1.5), Inches(0.38),
              f"{sc[key]:.1f} / {a['max_points']}", size=11, bold=True,
              color=color_rgb, font=HEADER, align=PP_ALIGN.RIGHT,
              anchor=MSO_ANCHOR.MIDDLE)
        # Bar
        bar_w = Inches(1.0 * sc[key] / a["max_points"]) if a["max_points"] > 0 else Inches(0.01)
        _rect(s2, Inches(11.5), row_y + Inches(0.14),
              Inches(1.0), Inches(0.1), fill=SLATE_200, corner=0.5)
        if bar_w > 0:
            _rect(s2, Inches(11.5), row_y + Inches(0.14), bar_w, Inches(0.1),
                  fill=color_rgb, corner=0.5)

    # Bottom: key findings
    _rect(s2, Inches(0.5), Inches(4.4), Inches(12.35), Inches(2.5),
          fill=INDIGO_LIGHT, corner=0.03)
    _text(s2, Inches(0.75), Inches(4.55), Inches(11), Inches(0.4),
          "主要所見", size=13, bold=True, color=SLATE_800, font=HEADER)
    strengths = project.get("findings", {}).get("strengths", [])
    weaknesses = project.get("findings", {}).get("weaknesses", [])

    _text(s2, Inches(0.75), Inches(5.0), Inches(5.8), Inches(0.35),
          "Strengths", size=11, bold=True, color=EMERALD, font=HEADER)
    sy = 5.35
    for i, s in enumerate(strengths[:3]):
        _text(s2, Inches(0.75), Inches(sy + i * 0.4), Inches(5.8), Inches(0.4),
              f"・{s}", size=11, color=SLATE_800, font=BODY)
    if not strengths:
        _text(s2, Inches(0.75), Inches(sy), Inches(5.8), Inches(0.4),
              "(アプリ内で入力してください)", size=10, color=SLATE_500,
              italic=True, font=BODY)

    _text(s2, Inches(7.0), Inches(5.0), Inches(5.8), Inches(0.35),
          "Weaknesses", size=11, bold=True, color=RED, font=HEADER)
    for i, w in enumerate(weaknesses[:3]):
        _text(s2, Inches(7.0), Inches(sy + i * 0.4), Inches(5.8), Inches(0.4),
              f"・{w}", size=11, color=SLATE_800, font=BODY)
    if not weaknesses:
        _text(s2, Inches(7.0), Inches(sy), Inches(5.8), Inches(0.4),
              "(アプリ内で入力してください)", size=10, color=SLATE_500,
              italic=True, font=BODY)

    _footer(s2, 2, TOTAL)

    # -------- Slide 3: Method --------
    s3 = prs.slides.add_slide(blank)
    _rect(s3, Inches(0), Inches(0), SW, SH, fill=SLATE_50)
    _title_bar(s3, "本診断の目的と方法", num=3)
    _text(s3, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
          "Purpose", size=13, bold=True, color=INDIGO, font=HEADER)
    _rect(s3, Inches(0.5), Inches(1.65), Inches(0.08), Inches(4.5), fill=INDIGO)
    _text(s3, Inches(0.75), Inches(1.65), Inches(5.8), Inches(4.5),
          "Webサイトリニューアル提案にあたり、現状のAIO/SEO観点での立ち位置を可視化し、"
          "リニューアルで強化すべき優先領域を特定することを目的としています。\n\n"
          "6領域・100点満点のスコアリングと、競合3-5社とのベンチマーク比較により、"
          "「なぜこの設計か」の根拠をデータで示します。",
          size=12, color=SLATE_800, font=BODY)

    _text(s3, Inches(7), Inches(1.2), Inches(6), Inches(0.4),
          "Overview", size=13, bold=True, color=INDIGO, font=HEADER)
    overview = [
        ("実施日", datetime.now().strftime("%Y/%m/%d")),
        ("対象URL", project.get("target_url", "")),
        ("業界カテゴリ", project.get("industry", "")),
        ("評価領域", "6領域 / 100点満点"),
        ("競合社", f"{len(project.get('competitors', []))}社"),
        ("担当", project.get("owner_email", "")),
    ]
    for i, (k, v) in enumerate(overview):
        row_y = Inches(1.65 + i * 0.55)
        fill = WHITE if i % 2 == 0 else SLATE_100
        _rect(s3, Inches(7), row_y, Inches(6), Inches(0.55), fill=fill)
        _text(s3, Inches(7.15), row_y, Inches(1.8), Inches(0.55),
              k, size=11, bold=True, color=SLATE_700, font=HEADER,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s3, Inches(9), row_y, Inches(3.9), Inches(0.55),
              str(v), size=11, color=SLATE_800, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s3, 3, TOTAL)

    # -------- Slide 4: 6 Areas --------
    s4 = prs.slides.add_slide(blank)
    _rect(s4, Inches(0), Inches(0), SW, SH, fill=SLATE_50)
    _title_bar(s4, "評価軸: 6領域 / 100点満点", num=4)
    col_x = [Inches(0.5), Inches(4.65), Inches(8.8)]
    row_y = [Inches(1.2), Inches(4.2)]
    for i, key in enumerate(AREA_ORDER):
        a = AREAS[key]
        x = col_x[i % 3]
        y = row_y[i // 3]
        c_main = RGBColor.from_string(a["color"].replace("#", ""))
        _rect(s4, x, y, Inches(4.05), Inches(2.8), fill=WHITE, corner=0.04)
        _rect(s4, x, y, Inches(0.15), Inches(2.8), fill=c_main)
        _text(s4, x + Inches(0.4), y + Inches(0.3), Inches(3.5), Inches(0.5),
              f"#{a['num']:02d}", size=14, bold=True, color=c_main, font=HEADER)
        _text(s4, x + Inches(0.4), y + Inches(0.8), Inches(3.5), Inches(0.6),
              a["name"], size=16, bold=True, color=SLATE_800, font=HEADER)
        _text(s4, x + Inches(0.4), y + Inches(1.4), Inches(3.5), Inches(0.5),
              a["desc"], size=11, color=SLATE_500, font=BODY)
        # Score
        _text(s4, x + Inches(0.4), y + Inches(2.0), Inches(3.5), Inches(0.6),
              f"{sc[key]:.1f}", size=28, bold=True, color=c_main, font=HEADER)
        _text(s4, x + Inches(2.3), y + Inches(2.15), Inches(1.5), Inches(0.5),
              f"/ {a['max_points']}", size=12, color=SLATE_500, font=BODY)
    _footer(s4, 4, TOTAL)

    # -------- Slide 5: Gap Analysis --------
    s5 = prs.slides.add_slide(blank)
    _rect(s5, Inches(0), Inches(0), SW, SH, fill=SLATE_50)
    _title_bar(s5, "ギャップ分析 / 強み・弱み", num=5)
    # Strengths
    _rect(s5, Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.6),
          fill=EMERALD_LIGHT, corner=0.04)
    _rect(s5, Inches(0.5), Inches(1.2), Inches(6.1), Inches(0.7),
          fill=EMERALD, corner=0.04)
    _text(s5, Inches(0.75), Inches(1.2), Inches(5.6), Inches(0.7),
          "STRENGTHS / 強み", size=14, bold=True, color=WHITE, font=HEADER,
          anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(strengths[:5]):
        _text(s5, Inches(0.8), Inches(2.0 + i * 0.9), Inches(5.5), Inches(0.8),
              f"・{item}", size=12, color=SLATE_800, font=BODY)
    if not strengths:
        _text(s5, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
              "アプリ内で所見を入力してください", size=11, color=SLATE_500,
              italic=True, font=BODY)

    # Weaknesses
    _rect(s5, Inches(6.75), Inches(1.2), Inches(6.1), Inches(5.6),
          fill=RED_LIGHT, corner=0.04)
    _rect(s5, Inches(6.75), Inches(1.2), Inches(6.1), Inches(0.7),
          fill=RED, corner=0.04)
    _text(s5, Inches(7.0), Inches(1.2), Inches(5.6), Inches(0.7),
          "WEAKNESSES / 弱み", size=14, bold=True, color=WHITE, font=HEADER,
          anchor=MSO_ANCHOR.MIDDLE)
    for i, item in enumerate(weaknesses[:5]):
        _text(s5, Inches(7.05), Inches(2.0 + i * 0.9), Inches(5.5), Inches(0.8),
              f"・{item}", size=12, color=SLATE_800, font=BODY)
    if not weaknesses:
        _text(s5, Inches(7.05), Inches(2.0), Inches(5.5), Inches(0.5),
              "アプリ内で所見を入力してください", size=11, color=SLATE_500,
              italic=True, font=BODY)
    _footer(s5, 5, TOTAL)

    # -------- Slide 6: Action Plan --------
    s6 = prs.slides.add_slide(blank)
    _rect(s6, Inches(0), Inches(0), SW, SH, fill=SLATE_50)
    _title_bar(s6, "改善アクション", num=6)
    actions = project.get("actions", [])
    if not actions:
        _text(s6, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
              "アプリ内で改善アクションを登録してください",
              size=13, italic=True, color=SLATE_500, font=BODY,
              align=PP_ALIGN.CENTER)
    else:
        # Header
        hdrs = [("No", 0.5), ("施策", 6.0), ("対象", 1.8), ("Impact", 1.3),
                ("Effort", 1.3), ("Period", 1.5)]
        x_acc = 0.5
        y = 1.2
        for h, w in hdrs:
            _rect(s6, Inches(x_acc), Inches(y), Inches(w), Inches(0.5), fill=SLATE_800)
            _text(s6, Inches(x_acc), Inches(y), Inches(w), Inches(0.5),
                  h, size=11, bold=True, color=WHITE, font=HEADER,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x_acc += w
        # Rows
        for i, a in enumerate(actions[:8]):
            ry = Inches(1.7 + i * 0.55)
            _rect(s6, Inches(0.5), ry, Inches(12.4), Inches(0.55),
                  fill=WHITE if i % 2 == 0 else SLATE_100)
            x_acc = 0.5
            vals = [
                str(i + 1), a.get("title", ""), a.get("area", ""),
                a.get("impact", ""), a.get("effort", ""), a.get("period", ""),
            ]
            for j, (val, (_, w)) in enumerate(zip(vals, hdrs)):
                _text(s6, Inches(x_acc + 0.1), ry, Inches(w), Inches(0.55),
                      str(val), size=10, color=SLATE_800, font=BODY,
                      align=PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
                x_acc += w
    _footer(s6, 6, TOTAL)

    # -------- Slide 7: Next Step --------
    s7 = prs.slides.add_slide(blank)
    _rect(s7, Inches(0), Inches(0), SW, SH, fill=SLATE_900)
    _rect(s7, Inches(0.5), Inches(0.5), Inches(0.2), Inches(0.6), fill=INDIGO)
    _text(s7, Inches(0.8), Inches(0.5), Inches(12), Inches(0.6),
          "NEXT STEP", size=14, bold=True, color=INDIGO, font=HEADER,
          anchor=MSO_ANCHOR.MIDDLE)
    _text(s7, Inches(0.5), Inches(1.2), Inches(12.35), Inches(1.0),
          "UI/UXリニューアルと並走した", size=28, bold=True, color=WHITE, font=HEADER)
    _text(s7, Inches(0.5), Inches(1.9), Inches(12.35), Inches(1.0),
          "AIO/SEO 強化プラン", size=36, bold=True, color=WHITE, font=HEADER)

    # Gp scope
    _rect(s7, Inches(0.5), Inches(3.2), Inches(6.1), Inches(3.4),
          fill=SLATE_800, corner=0.04)
    _text(s7, Inches(0.75), Inches(3.35), Inches(5.6), Inches(0.4),
          "Goodpatchがご提案する領域", size=14, bold=True, color=INDIGO, font=HEADER)
    scope = [
        "UI/UXリニューアル (主領域)", "AIO/SEO戦略設計 (本診断ベース)",
        "構造化データ設計 (CMS連携)", "E-E-A-T戦略ページ構築",
        "AI基盤構築 (パートナー連携)",
    ]
    for i, item in enumerate(scope):
        yy = Inches(3.85 + i * 0.5)
        _rect(s7, Inches(0.75), yy + Inches(0.17), Inches(0.12), Inches(0.12),
              fill=INDIGO, corner=0.5)
        _text(s7, Inches(1.0), yy, Inches(5.3), Inches(0.5),
              item, size=13, color=WHITE, font=BODY, anchor=MSO_ANCHOR.MIDDLE)

    # Agenda
    _rect(s7, Inches(6.75), Inches(3.2), Inches(6.1), Inches(3.4),
          fill=INDIGO, corner=0.04)
    _text(s7, Inches(7.0), Inches(3.35), Inches(5.6), Inches(0.4),
          "次回MTGのアジェンダ案", size=14, bold=True, color=WHITE, font=HEADER)
    agenda = [
        "1. 診断結果への Q&A (15分)",
        "2. 優先領域の合意形成 (20分)",
        "3. 提案スコープの握り (15分)",
        "4. スケジュール・体制確認 (10分)",
        "5. 次回までの宿題確認 (5分)",
    ]
    for i, item in enumerate(agenda):
        yy = Inches(3.85 + i * 0.5)
        _text(s7, Inches(7.0), yy, Inches(5.6), Inches(0.5),
              item, size=13, color=WHITE, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    _rect(s7, Inches(0), Inches(6.85), SW, Inches(0.65), fill=INDIGO)
    _text(s7, Inches(0.5), Inches(6.85), Inches(12.35), Inches(0.65),
          "本診断結果を起点に、御社のリニューアルをご一緒させてください。",
          size=14, bold=True, color=WHITE, font=HEADER,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=True)

    # -------- Slide 8: Thank You --------
    s8 = prs.slides.add_slide(blank)
    _rect(s8, Inches(0), Inches(0), SW, SH, fill=SLATE_900)
    _text(s8, Inches(10.0), Inches(0.2), Inches(3.2), Inches(3.2),
          f"{TOTAL:02d}", size=240, bold=True, color=SLATE_800, font=HEADER,
          align=PP_ALIGN.RIGHT)
    _rect(s8, Inches(0.75), Inches(2.5), Inches(0.15), Inches(1.2), fill=INDIGO)
    _text(s8, Inches(1.1), Inches(2.5), Inches(11), Inches(0.5),
          "Thank You", size=14, bold=True, color=INDIGO, font=HEADER)
    _text(s8, Inches(1.1), Inches(3.0), Inches(11), Inches(1.2),
          "ご清聴ありがとうございました", size=44, bold=True, color=WHITE, font=HEADER)
    _rect(s8, Inches(0.75), Inches(4.8), Inches(11.8), Inches(1.5),
          fill=SLATE_800, corner=0.03)
    _text(s8, Inches(1), Inches(4.95), Inches(4), Inches(0.4),
          "Contact", size=11, bold=True, color=INDIGO, font=HEADER)
    _text(s8, Inches(1), Inches(5.3), Inches(11), Inches(0.4),
          project.get("owner_email", "担当者"), size=14, bold=True, color=WHITE, font=HEADER)
    _text(s8, Inches(1), Inches(5.7), Inches(11), Inches(0.4),
          "Goodpatch Inc. / Market Design Div", size=12, color=SLATE_300, font=BODY)
    _rect(s8, Inches(0), Inches(7.0), SW, Inches(0.5), fill=SLATE_800)
    _text(s8, Inches(0.5), Inches(7.0), Inches(12.35), Inches(0.5),
          "AIO-Scope by Goodpatch",
          size=10, color=SLATE_300, font=BODY,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
